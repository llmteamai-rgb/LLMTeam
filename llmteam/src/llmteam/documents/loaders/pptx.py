"""
PowerPoint loader for .pptx files.

Uses python-pptx to load presentations.
"""

from pathlib import Path
from typing import Any, Dict, List, Union

from llmteam.documents.models import Document
from llmteam.documents.loaders.base import BaseLoader


class PPTXLoader(BaseLoader):
    """
    Loader for PowerPoint files (.pptx).

    Extracts text from slides, including titles, body text, and notes.

    Requires:
        pip install python-pptx
    """

    SUPPORTED_EXTENSIONS = {".pptx"}

    def __init__(
        self,
        include_notes: bool = True,
        one_doc_per_slide: bool = False,
    ):
        """
        Initialize the PowerPoint loader.

        Args:
            include_notes: Include slide notes in output.
            one_doc_per_slide: Create one Document per slide (vs one for whole presentation).
        """
        self.include_notes = include_notes
        self.one_doc_per_slide = one_doc_per_slide

    def load(self, source: Union[str, Path]) -> List[Document]:
        """
        Load a PowerPoint file.

        Args:
            source: Path to the .pptx file.

        Returns:
            List of Documents (one per slide or one for entire presentation).
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "PowerPoint support requires 'python-pptx'. "
                "Install with: pip install python-pptx"
            )

        path = Path(source)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        prs = Presentation(str(path))
        documents: List[Document] = []

        if self.one_doc_per_slide:
            for slide_num, slide in enumerate(prs.slides, start=1):
                content = self._extract_slide_content(slide)
                if content.strip():
                    metadata: Dict[str, Any] = {
                        "source": str(path),
                        "source_type": "powerpoint",
                        "file_name": path.name,
                        "slide_number": slide_num,
                        "total_slides": len(prs.slides),
                    }
                    documents.append(Document(content=content, metadata=metadata))
        else:
            # Combine all slides into one document
            all_content: List[str] = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                content = self._extract_slide_content(slide)
                if content.strip():
                    all_content.append(f"--- Slide {slide_num} ---\n{content}")

            if all_content:
                full_content = "\n\n".join(all_content)
                metadata = {
                    "source": str(path),
                    "source_type": "powerpoint",
                    "file_name": path.name,
                    "total_slides": len(prs.slides),
                }
                documents.append(Document(content=full_content, metadata=metadata))

        return documents

    def _extract_slide_content(self, slide) -> str:
        """Extract text content from a slide."""
        parts: List[str] = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())

        # Extract notes if enabled
        if self.include_notes and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_frame = notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                parts.append(f"[Notes: {notes_frame.text.strip()}]")

        return "\n".join(parts)

    def supports(self, path: Union[str, Path]) -> bool:
        """Check if this loader supports the given file."""
        return self._get_extension(path) in self.SUPPORTED_EXTENSIONS


__all__ = ["PPTXLoader"]
